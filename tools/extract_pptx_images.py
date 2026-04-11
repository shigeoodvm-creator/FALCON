import os
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx_path = "C:/FALCON/reference/FALCON画面.pptx"
output_dir = "C:/FALCON/docs/img"

os.makedirs(output_dir, exist_ok=True)

prs = Presentation(pptx_path)

results = []
img_counter = [0]  # リストでmutable参照


def extract_images_from_shapes(shapes, slide_num, slide_info):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            extract_images_from_shapes(shape.shapes, slide_num, slide_info)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_counter[0] += 1
            image = shape.image
            ext = image.ext
            filename = "slide{:02d}_img{:02d}.{}".format(slide_num, img_counter[0], ext)
            filepath = os.path.join(output_dir, filename)
            with open(filepath, "wb") as f:
                f.write(image.blob)
            slide_info["images"].append({
                "filename": filename,
                "size_bytes": len(image.blob),
                "content_type": image.content_type
            })
            print("  保存: {} ({:,} bytes)".format(filename, len(image.blob)))


for slide_num, slide in enumerate(prs.slides, start=1):
    # スライドのテキストを収集
    slide_texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_texts.append(text)

    slide_info = {
        "slide": slide_num,
        "texts": slide_texts,
        "images": []
    }

    print("\nスライド {}:".format(slide_num))
    if slide_texts:
        print("  テキスト: {}".format(" / ".join(slide_texts[:10])))
    else:
        print("  テキストなし")

    extract_images_from_shapes(slide.shapes, slide_num, slide_info)

    if not slide_info["images"]:
        print("  画像なし")

    results.append(slide_info)

print("\n=== 合計 {} 枚の画像を抽出しました ===".format(img_counter[0]))

# JSON形式でサマリーを保存
summary_path = os.path.join(output_dir, "extraction_summary.json")
with open(summary_path, "w", encoding="utf-8") as f:
    json.dump(results, f, ensure_ascii=False, indent=2)
print("サマリー保存: {}".format(summary_path))
